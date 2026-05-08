"""
generate_pptx.py — Copy a reference PPTX and update KPI text via find-replace.

Modes:
  --scan    Print text content of slide 1 and slide 2 from the reference file (no copy, no save).
  normal    Copy reference to output, apply replacements on ALL slides, save.

Usage (scan):
  python generate_pptx.py --reference "RG/.../ref.pptx" --scan

Usage (update):
  python generate_pptx.py \
    --reference "RG/2026/4. Abril/9. RG. 21.04.2026 Sistemas 03 - 2026.pptx" \
    --output    "RG/2026/5. Mayo/9. RG. 21.05.2026 Sistemas 04 - 2026.pptx"   \
    --find-replace "Marzo 2026" "Abril 2026" \
    --find-replace "21 de Abril, 2026" "21 de Mayo, 2026" \
    --find-replace "68" "50" \
    --find-replace "54" "38" \
    [--find-replace OLD NEW ...]
"""

import argparse
import os
import shutil
import sys

from pptx import Presentation


def para_full_text(para):
    return "".join(run.text for run in para.runs)


def replace_in_paragraph(para, old, new):
    """Replace old with new in a paragraph, preserving the format of the first matching run."""
    full = para_full_text(para)
    if old not in full:
        return False
    # Put the replaced text into the first run, wipe the rest
    new_full = full.replace(old, new)
    if para.runs:
        para.runs[0].text = new_full
        for run in para.runs[1:]:
            run.text = ""
    return True


def apply_replacements(slide, replacements):
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for old, new in replacements:
                if replace_in_paragraph(para, old, new):
                    count += 1
    return count


def print_slide_text(slide, slide_num, label=""):
    tag = f"Slide {slide_num}" + (f" — {label}" if label else "")
    print(f"\n{'='*60}")
    print(tag)
    print(f"{'='*60}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        lines = []
        for para in shape.text_frame.paragraphs:
            t = para_full_text(para)
            if t.strip():
                lines.append(t)
        if lines:
            print(f"  [{shape.name}]:")
            for line in lines:
                try:
                    print(f"    {line}")
                except UnicodeEncodeError:
                    print(f"    {line.encode('ascii', 'replace').decode()}")


def main():
    parser = argparse.ArgumentParser(description="Generate RG PPTX from reference")
    parser.add_argument("--reference", required=True, help="Path to reference PPTX")
    parser.add_argument("--output", help="Path for output PPTX (required unless --scan)")
    parser.add_argument("--scan", action="store_true",
                        help="Print slide 1+2 text from reference only, no file changes")
    parser.add_argument("--find-replace", nargs=2, action="append", metavar=("OLD", "NEW"),
                        default=[], help="Text replacement pair (repeatable)")
    args = parser.parse_args()

    if not os.path.exists(args.reference):
        print(f"ERROR: Reference file not found: {args.reference}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(args.reference)

    if args.scan:
        print(f"Reference: {args.reference}  ({len(prs.slides)} slides)")
        for i in range(min(4, len(prs.slides))):
            print_slide_text(prs.slides[i], i + 1)
        return

    if not args.output:
        print("ERROR: --output is required unless --scan is used.", file=sys.stderr)
        sys.exit(1)

    # Copy reference to output
    output_dir = os.path.dirname(args.output)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    shutil.copy2(args.reference, args.output)
    print(f"Copied reference -> {args.output}")

    # Re-open the copy for editing
    prs = Presentation(args.output)

    print("\n--- BEFORE ---")
    for i in range(min(4, len(prs.slides))):
        print_slide_text(prs.slides[i], i + 1)

    # Apply replacements on ALL slides
    total = 0
    for slide in prs.slides:
        total += apply_replacements(slide, args.find_replace)

    print("\n--- AFTER ---")
    for i in range(min(4, len(prs.slides))):
        print_slide_text(prs.slides[i], i + 1)

    prs.save(args.output)
    print(f"\nSaved: {args.output}")
    print(f"Total replacements made: {total}")

    if total == 0 and args.find_replace:
        print("\nWARNING: 0 replacements. Check that OLD strings match exactly (case-sensitive).")
    elif not args.find_replace:
        print("\nNo --find-replace pairs given. File copied with no text changes.")
        print("Run with --scan first to see current text, then re-run with --find-replace pairs.")


if __name__ == "__main__":
    main()
